# -*- coding: utf-8 -*-
"""
Created on Wed May 12 05:48:03 2021

@author: B57876
"""

# Create a .ppt to summarize results of a run of the microgrid(s)

# IMPORT
# standard
import os
import sys
from pptx import Presentation
import datetime
from datetime import timedelta
import numpy as np
import pandas as pd
from PIL import Image, ImageDraw, ImageFont
# perso
from ppt_utils import add_img_to_slide, init_img_plus_title_slide, add_linebreak_to_txt, set_txt_location, \
    suppress_unused_text_placeholders
from calc_output_metrics import check_if_unique_list
from plot import plot_mg_load_during_coord_method, plot_all_teams_mg_load_last_iter, \
            plot_per_actor_load_last_iter, plot_all_teams_two_metrics_tradeoff_last_iter, \
            plot_all_teams_score_traj, plot_agent_results_comparison

# Cf. https://python-pptx.readthedocs.io/en/latest/
# https://python-pptx.readthedocs.io/en/latest/user/quickstart.html

# Global params
IMG_FORMAT = "png"
IMG_SLIDE_LAYOUT_IDX = 1  # layout used for the slides with images
#    table_slide_layout_idx = 1 # and the one for slides with tables
LEFT_EMPTY_SPACE = 2  # in percentage of slide width
BOTTOM_EMPTY_SPACE = 5  # idem with height
FONT_STYLE = {"name": "Calibri", "size": 30, "bold": True, "italic": False}
# for team name style
TEAM_NAME_FONT_STYLE = {"name": "Calibri", "size": 15, "bold": True, "italic": False}
# and associated textbox
TEAM_NAME_TEXTBOX_STYLE = {"outline_color": "black", "outline_width": 2, "outline_style": "solid", "fill_color": "white"}
TEXT_VERTICAL_ALIGN = "middle"  # "top", "middle", "bottom"
DATE_FORMAT = "%Y-%m-%d %H:%M"
FILE_DATE_FORMAT = "%Y-%m-%d_%H%M"
# for the table not to be too width
SHORT_REGION_NAMES = {"grand_nord": "gd N", "grand_est": "gd E",
                      "grand_rhone": "gd Rhone", "bretagne": "bret.",
                      "grand_ouest": "gd O", "grand_sud_ouest": "gd SO",
                      "grande_ardeche": "gde ard.", "grand_sud_est": "gd SE"}
# parameters to locate text on the (French) regions map
# coordinates of the "centers" of the different regions
REGIONS_CENTER = {"grand_nord": (270, 70), "grand_est": (450, 170),
                  "grand_rhone": (425, 325), "bretagne": (50, 150),
                  "grand_ouest": (232, 265), "grand_sud_ouest": (230, 455),
                  "grande_ardeche": (455, 427), "grand_sud_est": (440, 520)}
REGIONS_TEXTBOX_MAX_WIDTH = {"grand_nord": 50, "grand_est": 120,
                             "grand_rhone": 86, "bretagne": 120,
                             "grand_ouest": 110, "grand_sud_ouest": 90,
                             "grande_ardeche": 75, "grand_sud_est": 90}


def get_teams_of_run(team_scores: dict) -> (list, int):
    """
    Get teams which take part of current run, and their number
    """

    team_names = list(team_scores)
    n_teams = len(team_names)

    return team_names, n_teams


def get_regions(team_scores: dict, team_names: list) -> list:
    """
    Get names of the regions in which current run has been made
    """

    regions = [list(team_scores[team]) for team in team_names]

    # check that all teams have the same regions scenarios
    if not check_if_unique_list(my_list_of_lists=regions):
        print("Different teams does not have the same region names... -> STOP")
        sys.exit(1)
    else:
        return regions[0]


def set_to_multiple_scenarios_format(dict_wo_scenarios: dict, fixed_scenario: tuple=(1,1,1,1)) -> dict:
    """
    Extend a dictionary to a new format in which the scenario indices of the different agent data is provided
    """

    dict_with_multiple_scenarios = {fixed_scenario[0]:
                                        {fixed_scenario[1]:
                                             {fixed_scenario[2]:
                                                  {fixed_scenario[3]: dict_wo_scenarios}
                                              }
                                         }
                                    }

    return dict_with_multiple_scenarios


class PptSynthesis():
    def __init__(self, result_dir: str, date_of_run: datetime.datetime, idx_run: int, optim_period: pd.date_range,
                 coord_method: str, regions_map_file: str):
        self.result_dir = os.path.join(result_dir, f'run_{date_of_run:%Y-%m-%d-%H%M}')
        os.makedirs(self.result_dir, exist_ok=True)
        self.date_of_run = date_of_run
        self.idx_run = idx_run
        self.optim_period = optim_period
        self.coord_method = coord_method
        self.regions_map_file = regions_map_file

    def create_title_slide(self):
        self.title_slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(self.title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"Summary of microgrid serious game run {self.idx_run}"
        subtitle.text = ("Generated on {:%s}" % DATE_FORMAT).format(self.date_of_run)

    def add_sized_img_to_slide(self, slide, img: Image, img_file: str, title_shape):
        add_img_to_slide(slide, img, img_file,
                         ((1 - 2 * LEFT_EMPTY_SPACE / 100) * self.prs.slide_width,
                          (1 - 2 * BOTTOM_EMPTY_SPACE / 100) * self.prs.slide_height \
                          - title_shape.height),
                         title_shape.height, LEFT_EMPTY_SPACE / 100 * self.prs.slide_width,
                         BOTTOM_EMPTY_SPACE / 100 * self.prs.slide_height)

    def create_summary_of_run_ppt(self, pv_prof: np.ndarray, load_profiles: dict, microgrid_prof: dict,
                                  microgrid_pmax: dict, per_actor_bills_internal: dict, cost_autonomy_tradeoff: dict,
                                  cost_co2emis_tradeoff: dict, team_scores: dict, best_teams_per_region: dict, scores_traj: dict):
        """
        Create a powerpoint to summarize the results of a given run of the microgrid
        serious game

        :param pv_prof: PV prod. profile
        :param load_profiles: dictionary with keys 1. Industrial Cons. scenario;
        2. Data Center scenario; 3. PV scenario; 4. EV scenario; 5. microgrid team name;
        6. Iteration; 7. actor type (N.B. charging_station_1, charging_station_2...
        if multiple charging stations in this microgrid) and values the associated
        load profile (kW)
        :param microgrid_prof: dict. with keys 1. Industrial Cons. scenario;
        2. Data Center scenario; 3. PV scenario; 4. EV scenario; 5. microgrid team name;
        6. Iteration and value the associated aggreg. microgrid load profile
        :param microgrid_pmax: dict. with keys 1. Industrial Cons. scenario;
        2. Data Center scenario; 3. PV scenario; 4. EV scenario; 5. microgrid team name;
        6. Iteration and value the associated microgrid pmax
        :param cost_autonomy_tradeoff: dict. with keys 1. the team names; 2. PV region
        names and values the associated (cost, autonomy score) aggreg. over the set
        of other scenarios
        :param per_actor_bills_internal: per actor INTERNAL bills, i.e. based on the MG price coordination
         signal. Same keys as preceding dict.
        :param cost_co2emis_tradeoff: idem with (cost, CO2 emissions) tradeoff
        :param team_scores: dict. with keys 1. team name; 2. region name and values
        the associated score of the current run
        :param best_teams_per_region: dict. with keys the names of the region simulated
        in current run and values the associated list of best team(s)
        :param scores_traj: dict. with keys 1. team name ; 2. run dates and associated
        values the score
        """

        region_coord_dyn_plot = None # if None the first region will be used

        # create an empty presentation
        self.prs = Presentation()

        # Create a title slide first
        self.create_title_slide()

        # DETAILED results
        # get team and region names
        team_names, n_teams = get_teams_of_run(team_scores=team_scores)
        regions = get_regions(team_scores=team_scores, team_names=team_names)

        # set region for coord dyn plot if not provided in parameter
        if region_coord_dyn_plot is None:
            region_coord_dyn_plot = regions[0]
        print(f"Region used for coordination dyn. plot: {region_coord_dyn_plot}")

        # TODO 2021-5-17 exclude teams with bug
        # slide to list these teams

        # 1) 1 slide per team with per iteration total microgrid load
        for i_team in range(n_teams):
            team = team_names[i_team]
            slide, shapes, title_shape = \
                init_img_plus_title_slide(self.prs, IMG_SLIDE_LAYOUT_IDX,
                                          f"Team {team} load DURING {self.coord_method} \n (PV) region: {region_coord_dyn_plot}",
                                          FONT_STYLE["name"], FONT_STYLE["size"], FONT_STYLE["bold"], FONT_STYLE["italic"],
                                          TEXT_VERTICAL_ALIGN)

            # plot and save
            current_dyn_mg_load_file = os.path.join(self.result_dir,
                                                    f"mg_load_during_dyn_team_{team}_{self.date_of_run.strftime(FILE_DATE_FORMAT)}.{IMG_FORMAT}")
            plot_mg_load_during_coord_method(microgrid_prof, region_coord_dyn_plot,
                                             team, current_dyn_mg_load_file.split(".")[0],
                                             self.optim_period)

            # open as Pillow Image
            dyn_mg_load_img = Image.open(current_dyn_mg_load_file)

            # add sized image to slide
            self.add_sized_img_to_slide(slide=slide, img=dyn_mg_load_img, img_file=current_dyn_mg_load_file,
                                        title_shape=title_shape)

            # suppress unused text placeholder (of index 1, 0 is used for the title)
            suppress_unused_text_placeholders(shapes)

        # 1) 1 slide per team with last iteration per-actor load
        # TODO add possibility to plot best iter
        for i_team in range(n_teams):
            team = team_names[i_team]
            slide, shapes, title_shape = \
                init_img_plus_title_slide(self.prs, IMG_SLIDE_LAYOUT_IDX,
                                          f"Team {team} ACTORS' load AT THE END of {self.coord_method} \n (PV) region: {region_coord_dyn_plot}",
                                          FONT_STYLE["name"], FONT_STYLE["size"], FONT_STYLE["bold"], FONT_STYLE["italic"],
                                          TEXT_VERTICAL_ALIGN)

            # plot and save
            per_actor_load_file = os.path.join(self.result_dir,
                                               f"per_actor_load_last_iter_team_{team}_{self.date_of_run.strftime(FILE_DATE_FORMAT)}.{IMG_FORMAT}")
            plot_per_actor_load_last_iter(load_profiles, pv_prof, region_coord_dyn_plot,
                                          team, per_actor_load_file.split(".")[0],
                                          self.optim_period)
            # open as Pillow Image
            per_actor_load_img = Image.open(per_actor_load_file)

            # add sized image to current slide
            self.add_sized_img_to_slide(slide=slide, img=per_actor_load_img, img_file=per_actor_load_file,
                                        title_shape=title_shape)

            # suppress unused text placeholder (of index 1, 0 is used for the title)
            suppress_unused_text_placeholders(shapes)

        # 1 slide with all teams microgrid load at the last iteration
        # TODO add possibility to plot best iter
        slide, shapes, title_shape = \
            init_img_plus_title_slide(self.prs, IMG_SLIDE_LAYOUT_IDX,
                                      f"All teams load AT THE END of {self.coord_method} \n (PV) region: {region_coord_dyn_plot}",
                                      FONT_STYLE["name"], FONT_STYLE["size"], FONT_STYLE["bold"], FONT_STYLE["italic"],
                                      TEXT_VERTICAL_ALIGN)

        # plot and save
        all_teams_mg_load_file = os.path.join(self.result_dir,
                                              f"all_teams_mg_load_last_iter_{self.date_of_run.strftime(FILE_DATE_FORMAT)}.{IMG_FORMAT}")
        plot_all_teams_mg_load_last_iter(microgrid_prof, microgrid_pmax, pv_prof, region_coord_dyn_plot,
                                         all_teams_mg_load_file.split(".")[0], self.optim_period)

        # open as Pillow Image
        all_teams_mg_load_img = Image.open(all_teams_mg_load_file)

        # add sized image to current slide
        self.add_sized_img_to_slide(slide=slide, img=all_teams_mg_load_img, img_file=all_teams_mg_load_file,
                                    title_shape=title_shape)

        # suppress unused text placeholder (of index 1, 0 is used for the title)
        suppress_unused_text_placeholders(shapes)

        # 1 slide with a comparison of all agents INTERNAL cost at last iteration of coord. method
        slide, shapes, title_shape = \
            init_img_plus_title_slide(self.prs, IMG_SLIDE_LAYOUT_IDX,
                                      f"All teams agents norm. cost AT THE END of {self.coord_method} \n (PV) region: {region_coord_dyn_plot}",
                                      FONT_STYLE["name"], FONT_STYLE["size"], FONT_STYLE["bold"], FONT_STYLE["italic"],
                                      TEXT_VERTICAL_ALIGN)

        # plot and save
        # consider first scenario for this plot
        scenario_plot = {"ic": 1, "dc": 1, "pv": region_coord_dyn_plot, "ev": 1}
        per_agent_cost_comparison_file = os.path.join(self.result_dir,
                                                      f"per_agent_cost_comparison_last_iter_{self.date_of_run.strftime(FILE_DATE_FORMAT)}.{IMG_FORMAT}")
        plot_agent_results_comparison(per_actor_bills_internal=per_actor_bills_internal, scenario_plot=scenario_plot,
                                      filename=per_agent_cost_comparison_file.split(".")[0], save_fig=True)

        # open as Pillow Image
        per_agent_cost_comparison_img = Image.open(per_agent_cost_comparison_file)

        # add sized image to current slide
        self.add_sized_img_to_slide(slide=slide, img=per_agent_cost_comparison_img,
                                    img_file=per_agent_cost_comparison_file, title_shape=title_shape)

        # suppress unused text placeholder (of index 1, 0 is used for the title)
        suppress_unused_text_placeholders(shapes)

        # 1 slide with scatter with 1 (eur, autonomy) point per team*PV region
        slide, shapes, title_shape = \
            init_img_plus_title_slide(self.prs, IMG_SLIDE_LAYOUT_IDX, f"All teams (cost, autonomy) tradeoff with {self.coord_method}",
                                      FONT_STYLE["name"], FONT_STYLE["size"], FONT_STYLE["bold"], FONT_STYLE["italic"],
                                      TEXT_VERTICAL_ALIGN)

        # plot and save
        all_teams_cost_auton_tradeoff_file = \
                os.path.join(self.result_dir, f"all_teams_cost_auton_tradeoff_last_iter_{self.date_of_run.strftime(FILE_DATE_FORMAT)}.{IMG_FORMAT}")

        metric_labels = {"cost": "Cost (eur)", "autonomy_score": "Autonomy score"}
        plot_all_teams_two_metrics_tradeoff_last_iter(two_metrics_tradeoff=cost_autonomy_tradeoff,
                                                      metric_1="cost", metric_2="autonomy_score", metric_labels=metric_labels,
                                                      filename=all_teams_cost_auton_tradeoff_file.split(".")[0])

        # open as Pillow Image
        all_teams_cost_auton_tradeoff_img = Image.open(all_teams_cost_auton_tradeoff_file)

        # add sized image to current slide
        self.add_sized_img_to_slide(slide=slide, img=all_teams_cost_auton_tradeoff_img,
                                    img_file=all_teams_cost_auton_tradeoff_file, title_shape=title_shape)

        # suppress unused text placeholder (of index 1, 0 is used for the title)
        suppress_unused_text_placeholders(shapes)

        # 1 slide with scatter with 1 (eur, CO2 emissions) point per team*PV region
        slide, shapes, title_shape = \
            init_img_plus_title_slide(self.prs, IMG_SLIDE_LAYOUT_IDX, f"All teams (cost, CO2 emissions) tradeoff with {self.coord_method}",
                                      FONT_STYLE["name"], FONT_STYLE["size"], FONT_STYLE["bold"], FONT_STYLE["italic"],
                                      TEXT_VERTICAL_ALIGN)

        # plot and save
        all_teams_cost_co2emis_tradeoff_file = \
                os.path.join(self.result_dir, f"all_teams_cost_co2emis_tradeoff_last_iter_{self.date_of_run.strftime(FILE_DATE_FORMAT)}.{IMG_FORMAT}")

        metric_labels = {"cost": "Cost (eur)", "co2_emis": "CO2 emissions (gCO2eq.)"}
        plot_all_teams_two_metrics_tradeoff_last_iter(two_metrics_tradeoff=cost_co2emis_tradeoff,
                                                      metric_1="cost", metric_2="co2_emis", metric_labels=metric_labels,
                                                      filename=all_teams_cost_co2emis_tradeoff_file.split(".")[0])

        # open as Pillow Image
        all_teams_cost_co2emis_tradeoff_img = Image.open(all_teams_cost_co2emis_tradeoff_file)

        # add sized image to current slide
        self.add_sized_img_to_slide(slide=slide, img=all_teams_cost_co2emis_tradeoff_img,
                                    img_file=all_teams_cost_co2emis_tradeoff_file, title_shape=title_shape)

        # suppress unused text placeholder (of index 1, 0 is used for the title)
        suppress_unused_text_placeholders(shapes)

        # GLOBAL results
        # 1) best team per region
        # create image
        regions_map_img = self.create_best_team_per_region_img(best_teams_per_region)

        # init. slide with title
        slide, shapes, title_shape = \
            init_img_plus_title_slide(self.prs, IMG_SLIDE_LAYOUT_IDX,
                                      "Best team(s) per region (PV prof.)", FONT_STYLE["name"], FONT_STYLE["size"],
                                      FONT_STYLE["bold"], FONT_STYLE["italic"], TEXT_VERTICAL_ALIGN)

        # add sized image to current slide
        img_file = os.path.join(self.result_dir, f"best_team_per_region_{self.date_of_run.strftime(FILE_DATE_FORMAT)}.{IMG_FORMAT}")
        self.add_sized_img_to_slide(slide=slide, img=regions_map_img, img_file=img_file, title_shape=title_shape)

        # suppress unused text placeholder (of index 1, 0 is used for the title)
        suppress_unused_text_placeholders(shapes)

        # 2) table with scores per region for all teams (team names in lines, region
        # names in col. and associated score in each cell)
        # TODO put the table in the right placeholder...
        # TODO adapt location and size automatically
        # ALternative for now...
        slide = self.prs.slides.add_slide(self.title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Scores per team name*region (PV prof.)"
        subtitle.text = f"See file: aggreg_per_region_res_run_{self.date_of_run.strftime(FILE_DATE_FORMAT)}.csv"

    #    loc_x_table = 600
    #    loc_y_table = 500
    #    table_width = 2000
    #    table_height = 600
    #    # init. slide with title
    #    slide, shapes, title_shape = \
    #        init_img_plus_title_slide(prs, table_slide_layout_idx,
    #                                  "Scores per team name*region (PV prof.)", font_name,
    #                                  font_size, font_bold, font_italic, text_vertical_align)
    #    # get team and region names based on team_scores
    #    team_names = list(team_scores)
    #    n_teams = len(team_names)
    #    regions = [list(team_scores[team]) for team in team_names]
    #
    #    # check that all teams have the same regions scenarios
    #    if not check_if_unique_list(regions):
    #        print("Different teams does not have the same region names... -> STOP")
    #        sys.exit(1)
    #    else:
    #        regions = regions[0]
    #
    #    n_regions = len(regions)

        # initialize table
    #    table = slide.placeholders[1].insert_table(len(team_names)+1, len(regions)+1)
    #    table = slide.shapes.add_table(len(team_names)+1, len(regions)+1, loc_x_table,
    #                                   loc_y_table, table_width, table_height).table
        # fill each cell by looping over all of them
        # title one
    #    table.cell(0, 0).text = "Team/region names"
    #    # names of the teams in lines
    #    for i_team in range(n_teams):
    #        table.cell(i_team+1, 0).text = team_names[i_team]
    #    # idem for regions in columns
    #    for i_region in range(n_regions):
    #        table.cell(0, i_region+1).text = regions[i_region] if not \
    #        regions[i_region] in short_region_names else short_region_names[region[i_region]]
    #    # loop over (team, region) to fullfill the center cells
    #    for i_team in range(n_teams):
    #        for i_region in range(n_regions):
    #            table.cell(i_team+1, i_region+1).text = \
    #                       "%.3f" % team_scores[team_names[i_team]][regions[i_region]]

        # improvement slide based on multiple aggreg_results.csv files
        # first check if there is more than one run available
        n_run_dates = len(scores_traj[list(scores_traj)[0]])
        if n_run_dates == 0:
            print("No slide for score improvement, because unique run available for now")
        else:
            slide, shapes, title_shape = \
                init_img_plus_title_slide(self.prs, IMG_SLIDE_LAYOUT_IDX, f"All teams score traj. with {self.coord_method}",
                                          FONT_STYLE["name"], FONT_STYLE["size"], FONT_STYLE["bold"], FONT_STYLE["italic"],
                                          TEXT_VERTICAL_ALIGN)

            # plot and save
            all_teams_score_traj_file = os.path.join(self.result_dir,
                                                     f"all_teams_score_traj_to_run_{self.date_of_run.strftime(FILE_DATE_FORMAT)}.{IMG_FORMAT}")

            plot_all_teams_score_traj(scores_traj, all_teams_score_traj_file.split(".")[0])

            # open as Pillow Image
            all_teams_score_traj_img = Image.open(all_teams_score_traj_file)

            # add sized image to current slide
            self.add_sized_img_to_slide(slide=slide, img=all_teams_score_traj_img, img_file=all_teams_score_traj_file,
                                        title_shape=title_shape)

            # suppress unused text placeholder (of index 1, 0 is used for the title)
            suppress_unused_text_placeholders(shapes)

        # save the pptx JDP presentation
        self.prs.save(os.path.join(self.result_dir, f"run_summary_{self.date_of_run.strftime(FILE_DATE_FORMAT)}.pptx"))

    def create_best_team_per_region_img(self, best_teams_per_region: dict) -> Image:
        """
        Create image giving the best team for each of the regions

        :param best_teams_per_region: dict. with keys the names of the region simulated
        in current run and values the associated list of best team(s)
        """

        img_margin_x = 5
        img_margin_y = 5
        textbox_margin_x = 10 # margin box rounding text
        textbox_margin_y = 10
        # TODO logo of given srious gamme (Ponts + 2021?)
    #    logo_relative_size = 10 # in proportion of photo image size
    #    logo_margin_x = 5
    #    logo_margin_y = 5
    #    logo_background_color = ""
    #    logo_opacity_level = 0

        # open Pillow Image
        regions_map_img = Image.open(self.regions_map_file)
        draw_regions_map_img = ImageDraw.Draw(regions_map_img)

        for region in best_teams_per_region:
            # Join text if multiple teams with same score
            best_teams_txt = "/".join(best_teams_per_region[region])
            # get text width and height
            current_font = ImageFont.load_default()
            #truetype("%s%s.ttf" % (TEAM_NAME_FONT_STYLE["name"].lower(),
            #                                  "b" if TEAM_NAME_FONT_STYLE["bold"] else "i" \
            #                                  if TEAM_NAME_FONT_STYLE["italic"] else ""),
            #                                  TEAM_NAME_FONT_STYLE["size"])

            # add linebreak in text to fit with max. textbox width
            best_teams_txt = add_linebreak_to_txt(best_teams_txt, draw_regions_map_img,
                                                  current_font, REGIONS_TEXTBOX_MAX_WIDTH[region])
            textwidth, textheight = draw_regions_map_img.textsize(best_teams_txt,
                                                                  current_font)
            # set text location
            loc_x, loc_y = set_txt_location(REGIONS_CENTER[region], regions_map_img.width,
                                            regions_map_img.height, textwidth,
                                            textheight, textbox_margin_x+img_margin_x,
                                            textbox_margin_y+img_margin_y)

            # textbox rounded the team name(s)
            draw_regions_map_img \
                .rectangle((loc_x-textbox_margin_x, loc_y-textbox_margin_y,
                            loc_x+textwidth+textbox_margin_x,
                            loc_y+textheight+textbox_margin_y),
                            outline=TEAM_NAME_TEXTBOX_STYLE["outline_color"],
                            width=TEAM_NAME_TEXTBOX_STYLE["outline_width"],
                            fill=TEAM_NAME_TEXTBOX_STYLE["fill_color"])

            # Add text
            draw_regions_map_img.text((loc_x, loc_y), best_teams_txt, fill=(0,0,0),
                                      font=current_font)

        # save image with best team name(s)
        regions_map_img \
            .save(os.path.join(self.result_dir,
                               f"best_team_per_region_{self.date_of_run.strftime('%Y-%m-%d_%H%M')}.{IMG_FORMAT}")
                  )

        return regions_map_img


def create_current_run_dir(current_dir: str, date_of_run: datetime.datetime):
    """
    Create a directory where all results of current run will be saved
    
    :param current_dir: directory in which photos are stored
    :param date_of_run: date of current run
    :return: returns the full path of current run directory
    """
    
    current_run_dir = os.path.join(current_dir,
                                   "run_%s" % date_of_run.strftime("%Y-%m-%d_%H%M"))
    
    if os.path.exists(current_run_dir):
        print("Result dir. for run %s already exists -> it is emptied before processing current outputs" \
              % date_of_run.strftime("%Y-%m-%d_%H%M"))
        all_files_tb_suppr = [os.path.join(current_run_dir, elt_file) \
                            for elt_file in os.listdir(current_run_dir) \
                            if (elt_file.endswith(".png") or elt_file.endswith(".csv"))]
        
        for elt_file in all_files_tb_suppr:
            os.remove(elt_file)
    
    # create backup photos dir
    else:
        os.mkdir(current_run_dir)
    
    return current_run_dir

if __name__ =="__main__":
    current_dir = os.getcwd()
    date_of_run = datetime.datetime.now()
    idx_run = 1
    
    # create directory for current run
    result_dir = create_current_run_dir(current_dir, date_of_run)

    delta_t_s = 1800 # time-slot duration (s)
    n_ts = 48
    start_optim_period = datetime.datetime(2018,1,1)
    optim_period = pd.date_range(start=start_optim_period,
                                 end=start_optim_period+timedelta(hours=24),
                                 freq="%is" % delta_t_s)[:-1]
    coord_method = "price-coord. dyn."
    regions_map_file = os.path.join(current_dir, "images", "pv_regions_no_names.png")

    ppt_synthesis = PptSynthesis(result_dir=result_dir, date_of_run=date_of_run, idx_run=idx_run,
                                 optim_period=optim_period, coord_method=coord_method,
                                 regions_map_file=regions_map_file)
    best_teams_per_region_1 = {"grand_nord": ["team 1"], "grand_est": ["big team"], "grand_rhone": ["qui ne saute pas"],
                               "bretagne": ["bonnets rouges et blanc bonnet"], "grand_ouest": ["a la masse"],
                               "grand_sud_ouest": ["la garonne est viola"], "grande_ardeche": ["ca canyone pas mal"],
                               "grand_sud_est": ["le soleil c'est la vie"]
                            }

    team_scores = {"team_1": {"grand_nord": 12, "grand_est": 57, "grand_rhone": 69,
                               "bretagne": 35, "grand_ouest": 47, "grand_sud_ouest": 33,
                               "grande_ardeche": 72, "grand_sud_est": 22},
                   "team_2": {"grand_nord": 15, "grand_est": 52, "grand_rhone": 61,
                              "bretagne": 32, "grand_ouest": 41, "grand_sud_ouest": 31,
                              "grande_ardeche": 71, "grand_sud_est": 22},
                   "team_3": {"grand_nord": 9, "grand_est": 57, "grand_rhone": 75,
                              "bretagne": 30, "grand_ouest": 40, "grand_sud_ouest": 29,
                              "grande_ardeche": 56, "grand_sud_est": 18}
                  }

    # Collective metrics calculation
    load_profiles = {"team_1":
                          {1:
                               {"solar_farm": np.random.rand(n_ts),
                                "industrial_consumer": np.random.rand(n_ts),
                                "charging_station": np.random.rand(n_ts),
                                "data_center": np.random.rand(n_ts)},
                           2:
                               {"solar_farm": np.random.rand(n_ts),
                                "industrial_consumer": np.random.rand(n_ts),
                                "charging_station": np.random.rand(n_ts),
                                "data_center": np.random.rand(n_ts)},
                           3:
                               {"solar_farm": np.random.rand(n_ts),
                                "industrial_consumer": np.random.rand(n_ts),
                                "charging_station": np.random.rand(n_ts),
                                "data_center": np.random.rand(n_ts)}
                           },
                     "team_2":
                         {1:
                              {"solar_farm": np.random.rand(n_ts),
                               "industrial_consumer": np.random.rand(n_ts),
                               "charging_station": np.random.rand(n_ts),
                               "data_center": np.random.rand(n_ts)},
                          2:
                              {"solar_farm": np.random.rand(n_ts),
                               "industrial_consumer": np.random.rand(n_ts),
                               "charging_station": np.random.rand(n_ts),
                               "data_center": np.random.rand(n_ts)},
                          3:
                              {"solar_farm": np.random.rand(n_ts),
                               "industrial_consumer": np.random.rand(n_ts),
                               "charging_station": np.random.rand(n_ts),
                               "data_center": np.random.rand(n_ts)},
                          4:
                              {"solar_farm": np.random.rand(n_ts),
                               "industrial_consumer": np.random.rand(n_ts),
                               "charging_station": np.random.rand(n_ts),
                               "data_center": np.random.rand(n_ts)}
                          },
                     "team_3":
                         {1:
                              {"solar_farm": np.random.rand(n_ts),
                               "industrial_consumer": np.random.rand(n_ts),
                               "charging_station": np.random.rand(n_ts),
                               "data_center": np.random.rand(n_ts)},
                          2:
                              {"solar_farm": np.random.rand(n_ts),
                               "industrial_consumer": np.random.rand(n_ts),
                               "charging_station": np.random.rand(n_ts),
                               "data_center": np.random.rand(n_ts)},
                          3:
                              {"solar_farm": np.random.rand(n_ts),
                               "industrial_consumer": np.random.rand(n_ts),
                               "charging_station": np.random.rand(n_ts),
                               "data_center": np.random.rand(n_ts)},
                          4:
                              {"solar_farm": np.random.rand(n_ts),
                               "industrial_consumer": np.random.rand(n_ts),
                               "charging_station": np.random.rand(n_ts),
                               "data_center": np.random.rand(n_ts)}
                          }
                     }

    purchase_price = 0.10 + 0.1 * np.random.rand(n_ts)
    sale_price = 0.05 + 0.1 * np.random.rand(n_ts)
    delta_t_s = 1800
    contracted_p_tariffs = {6: 123.6, 9: 151.32, 12: 177.24, 15: 201.36,
                            18: 223.68, 24: 274.68, 30: 299.52, 36: 337.56}

    from calc_output_metrics import calc_per_actor_bills, calc_microgrid_collective_metrics, \
       calc_cost_autonomy_tradeoff_last_iter, get_best_team_per_region, \
           save_all_metrics_to_csv, save_per_region_score_to_csv, get_improvement_traj

    # calculate per-actor bill
    per_actor_bills = calc_per_actor_bills(load_profiles, purchase_price,
                                           sale_price, delta_t_s)

    # and microgrid collective metrics
    microgrid_prof, microgrid_pmax, collective_metrics = \
            calc_microgrid_collective_metrics(load_profiles, contracted_p_tariffs,
                                              delta_t_s)

    # and finally cost-autonomy tradeoff
    cost_autonomy_tradeoff = \
        calc_cost_autonomy_tradeoff_last_iter(per_actor_bills, collective_metrics)

    # Get best team per region
    coll_metrics_weights = {"pmax_cost": 1/365, "autonomy_score": 1,
                            "mg_transfo_aging": 0, "n_disj": 0}
    team_scores, best_teams_per_region, coll_metrics_names = \
            get_best_team_per_region(per_actor_bills, collective_metrics,
                                     coll_metrics_weights)

    # save detailed results to a .csv file
    metrics_not_saved = ["mg_transfo_aging", "n_disj"]
    save_all_metrics_to_csv(per_actor_bills, collective_metrics, coll_metrics_names,
                            coll_metrics_weights, metrics_not_saved, result_dir,
                            date_of_run)

    # and aggreg. per region .csv file
    save_per_region_score_to_csv(team_scores, result_dir, date_of_run)

    # get "improvement trajectory"
    list_of_run_dates = [datetime.datetime.strptime(elt[4:], "%Y-%m-%d_%H%M") \
                         for elt in os.listdir(current_dir) \
                         if (os.path.isdir(elt) and elt.startswith("run_"))]
    scores_traj = get_improvement_traj(current_dir, list_of_run_dates,
                                       list(team_scores))

    pv_prof = 5 * np.random.rand(n_ts)

    ppt_synthesis.create_summary_of_run_ppt(pv_prof=pv_prof, load_profiles=load_profiles, microgrid_prof=microgrid_prof,
                                            microgrid_pmax=microgrid_pmax, cost_autonomy_tradeoff=cost_autonomy_tradeoff,
                                            team_scores=team_scores, best_teams_per_region=best_teams_per_region_1,
                                            scores_traj=scores_traj)
