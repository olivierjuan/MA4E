import sys
import numpy as np
import pptx
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR
from PIL import Image, ImageDraw, ImageFont


def set_txt_location(txt_center_ideal_loc: tuple, img_width: (float, int),
                     img_height: (float, int), textwidth: (float, int),
                     textheight: (float, int), margin_x: (float, int),
                     margin_y: (float, int)) -> tuple:
    """
    Set text location

    :param txt_center_ideal_loc: (x,y) ideal location for the center of the
    textbox
    :param img_width: image width on which the text will be written
    :param img_height: idem height
    :param textwidth: text width
    :param textheight: text height
    :param margin_x: x-axis margin to be taken between text and the image limits
    :param margin_y: idem for y-axis    """

    # x location
    if textwidth + 2 * margin_x > img_width:
        print("TextWIDTH pb: the text to be added is bigger than image!")
        loc_x = txt_center_ideal_loc[0] - textwidth / 2
    else:
        # check if text width is in conflict with left/right limits of the image
        loc_x = txt_center_ideal_loc[0] - textwidth / 2 if \
            (txt_center_ideal_loc[0] - textwidth / 2 - margin_x >= 0 and
             txt_center_ideal_loc[0] + textwidth / 2 + margin_x <= img_width) \
            else margin_x if txt_center_ideal_loc[0] - textwidth / 2 - margin_x < 0 \
            else img_width - textwidth - margin_x

    # y location
    if textheight + 2 * margin_y > img_height:
        print("TextHEIGHT pb: the text to be added is bigger than image!")
        loc_y = txt_center_ideal_loc[1] - textheight / 2
    else:
        # check if text height conflicting with top/bottom img limits
        loc_y = txt_center_ideal_loc[1] - textheight / 2 if \
            (txt_center_ideal_loc[1] - textheight / 2 - margin_y >= 0 and
             txt_center_ideal_loc[1] + textheight / 2 + margin_y <= img_height) \
            else margin_y if txt_center_ideal_loc[1] - textheight / 2 - margin_y < 0 \
            else img_height - textheight - margin_y

    return loc_x, loc_y


def add_linebreak_to_txt(my_txt: str, img_draw: ImageDraw, text_font: ImageFont,
                         max_width: int) -> str:
    """
    Add linebreak in text to fit with a max. width limit
    """

    textwidth, textheight = img_draw.textsize(my_txt, text_font)

    if textwidth <= max_width:
        return my_txt
    else:
        # get number of characters which correspond to a line
        line_len = int(np.floor(max_width / textwidth * len(my_txt)))

        my_txt_with_linebreaks = ""
        i = 0
        while i < len(my_txt):
            current_line = my_txt[i:i + line_len].lstrip()
            if len(current_line) < line_len:
                delta_fullfill = line_len - len(current_line)
                current_line += my_txt[i + line_len:i + line_len + delta_fullfill]
            else:
                delta_fullfill = 0

            my_txt_with_linebreaks += "%s\n" % current_line
            i += line_len + delta_fullfill

        # suppress last linebreak
        return my_txt_with_linebreaks[:-1]


def add_img_to_slide(slide: pptx.slide.Slide, img_to_add: Image, img_file: str,
                     img_box: tuple, title_height: float, left_space: float,
                     bottom_space: float):
    """
    Add an image to a simple slide with a title and an image

    :param slide: slide in which image has to be added
    :param img_to_add: Pillow Image to be added
    :param img_file: full path to the file of this image
    :param img_box: image box dimensions (width, height)
    :param title_height: height of the title above the image in slide
    :param left_space: space to the left of the img_box in slide
    :param bottom_space: space at the bottom of the img_box (and above) in slide
    """

    img_resize = resize_img_in_box(img_to_add, img_box)

    # set location based on img new dimensions
    left_img = left_space + (img_box[0] - img_resize[0]) / 2
    top_img = title_height + bottom_space + (img_box[1] - img_resize[1]) / 2
    # add this image
    slide.shapes.add_picture(img_file, left_img, top_img, height=img_resize[1])


def init_img_plus_title_slide(prs: Presentation, layout_idx: int, title_text: str, font_name: str, font_size: int,
                              font_bold: bool, font_italic: bool, text_vertical_align: str):
    """
    Initialize ppt slide with title and an image
    """

    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = title_text
    # set title style
    set_text_style(title_shape.text_frame.paragraphs[0], font_name, font_size,
                   font_bold, font_italic, text_vertical_align)

    return slide, shapes, title_shape


def set_text_style(text_shape, font_name: str = "Calibri", font_size: int = 18,
                   font_bold: bool = True, font_italic: bool = False,
                   text_vertical_align: str = "center"):
    """
    Set text style of a given text shape
    """
    font = text_shape.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = font_bold
    font.italic = font_italic  # cause value to be inherited from theme
    if text_vertical_align == "top":
        text_shape.vertical_anchor = MSO_ANCHOR.TOP
    elif text_vertical_align == "middle":
        text_shape.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif text_vertical_align == "bottom":
        text_shape.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        # ERROR
        print("Unknown text vertical alignment %s -> STOP" % text_vertical_align)
        sys.exit(1)


def resize_img_in_box(img: Image, img_box_size: tuple) -> tuple:
    """
    Resize image based on the size of an associated box

    :param img: Pillow image
    :param img_box_size: (width, height) of the associated box
    :return: returns the tuple (width, height) with updated image dimensions
    """

    # sizing operation limited by height
    if img_box_size[0] / img.width >= img_box_size[1] / img.height:
        return img.width * img_box_size[1] / img.height, img_box_size[1]
    else:
        return img_box_size[0], img.height * img_box_size[0] / img.width


def suppress_unused_text_placeholders(shapes):
    sp = shapes[1].element
    sp.getparent().remove(sp)
